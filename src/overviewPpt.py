import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def findStartCell(table):
    tRows = len(table.rows)
    tCols = len(table.columns)   
    for i in range(0,tRows):
        for j in range(0, tCols):
            if table.cell(i,j).text_frame.text == 'Start':
                return i, j

def fillTable(dataframe, table, startRow, startCol):
    for i in range(0, dataframe.shape[0]):
        for j in range(0, dataframe.shape[1]):

            formattedDataFrame = dataframe.apply(lambda x: x.apply(lambda y: '{:,}'.format(y) if type(y) == int else y))
            
            formattedDataFrame = formattedDataFrame.astype(str).replace({
                'nan': '',
                'nan%': '***',
                'inf': '***',
                'inf%': '***',
                # 'None': '***'
            })

            value = formattedDataFrame.iloc[i,j]
            table.cell(i+startRow, j+startCol).text_frame.text = str(value)

            for paragraph in table.cell(i+startRow, j+startCol).text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(7.5)

def getDF(list, name):
    for df in list:
        if df.index.name == name:
            return df
    return None

def summaryPage(presentation, dataFrames, dateStart, dateEnd):
    overview = presentation.slides[0]
    for shape in overview.shapes:
        if shape.has_text_frame:
            if "PERIOD COVERING" in shape.text:
                shape.text = f"PERIOD COVERING: {dateStart.strftime('%m/%d/%Y')} - {dateEnd.strftime('%m/%d/%Y')}"
                shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.paragraphs[0].font.size = Pt(16)

        if shape.has_table:
            table = shape.table
            tableName = table.cell(0,0).text_frame.text # table name is stored in the first cell and column as of now.
            # if this changes, need to figure out what cell it is in and change the indices to look for the name
            rowLength = len(table.rows)
            colLength = len(table.columns)
            startrow, startcol = findStartCell(table)

            match tableName:
                case _ if "FORCE LEVELS" in tableName:
                    forceLevels = getDF(dataFrames, 'FORCE LEVELS')
                    fillTable(forceLevels, table, startrow, startcol)

                case _ if "LEVEL 2 INCIDENTS" in tableName:
                    level2 = getDF(dataFrames, 'LEVEL 2')
                    fillTable(level2, table, startrow, startcol)

                    table.cell(rowLength-1, colLength-1).text_frame.text = ''
                    table.cell(rowLength-1, colLength-2).text_frame.text = ''

                case _ if "FORCE USED BY MOS" in tableName:
                    forceMos = getDF(dataFrames, 'FORCE USED BY MOS')
                    fillTable(forceMos, table, startrow, startcol)

                    table.cell(rowLength-1, colLength-1).text_frame.text = ''
                    table.cell(rowLength-1, colLength-2).text_frame.text = ''

                case _ if "PLATOON" in tableName:
                    platoon = getDF(dataFrames, 'PLATOON')
                    fillTable(platoon, table, startrow, startcol)

                case _ if "VELOCITY" in tableName:
                    injMos = getDF(dataFrames, 'VELOCITY')
                    fillTable(injMos, table, startrow, startcol)
                    
                case _ if "TRIs WITH ARREST" in tableName:
                    triArrest = getDF(dataFrames, 'TRIs WITH ARREST')
                    triArrestsNoTotal = (triArrest.T[['TRIs w/ Arrests', 'TRIs w/ Arrests where MOS used force', 'Arrests', '% of Total Arrests']]).T
                    fillTable(triArrestsNoTotal, table, startrow, startcol)

                    table.cell(rowLength-1, colLength-1).text_frame.text = ''
                    table.cell(rowLength-1, colLength-2).text_frame.text = ''

                case _ if "INJURED MOS" in tableName:
                    injMos = getDF(dataFrames, 'INJURED MOS')
                    fillTable(injMos, table, startrow, startcol)

                    table.cell(rowLength-1, colLength-1).text_frame.text = ''
                    table.cell(rowLength-1, colLength-2).text_frame.text = ''

def breakdownPages(presentation, dataFrames, dateStart, dateEnd):
    for slide in [presentation.slides[1], presentation.slides[2], presentation.slides[3]]:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "PERIOD COVERING" in shape.text:
                    shape.text = f"PERIOD COVERING: {dateStart.strftime('%m/%d/%Y')} - {dateEnd.strftime('%m/%d/%Y')}"
                    shape.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].font.size = Pt(16)

            if shape.has_table:
                table = shape.table
                tableName = table.cell(0,0).text_frame.text
                startrow, startcol = findStartCell(table)

                match tableName:
                    case "FORCE LEVELS PERIOD":
                        forceLevelsPeriod = getDF(dataFrames, 'FORCE LEVELS PERIOD')
                        fillTable(forceLevelsPeriod, table, startrow, startcol)
                    case "FORCE LEVELS PERIOD %":
                        forceLevelsPeriodPercent = getDF(dataFrames, 'FORCE LEVELS PERIOD %')
                        fillTable(forceLevelsPeriodPercent, table, startrow, startcol)
                    case "FORCE LEVELS YTD":
                        forceLevelsYTD = getDF(dataFrames, 'FORCE LEVELS YTD')
                        fillTable(forceLevelsYTD, table, startrow, startcol)
                    case "SUBJECT INJURY LEVELS":
                        subInjLevels = getDF(dataFrames, 'SUBJECT INJURY LEVELS')
                        fillTable(subInjLevels, table, startrow, startcol)
                    case "INJURED MOS":
                        injMos = getDF(dataFrames, 'INJURED MOS')
                        fillTable(injMos, table, startrow, startcol)
                    case "REASON FORCE USED BY MOS":
                        reasonForce = getDF(dataFrames, 'REASON FORCE')
                        fillTable(reasonForce, table, startrow, startcol)
                    case "TRIs WITH ARRESTS":
                        triArrests = getDF(dataFrames, 'TRIs WITH ARRESTS')
                        triArrestsNoTotal = (triArrests.T[['TRIs w/ Arrests', 'TRIs w/ Arrests where MOS used force', 'Arrests', '% of Total Arrests']]).T
                        fillTable(triArrestsNoTotal, table, startrow, startcol)
                    case "TYPE OF PHYSICAL FORCE USED BY MOS":
                        typeForce = getDF(dataFrames, 'TYPE OF PHYSICAL FORCE')
                        fillTable(typeForce, table, startrow, startcol)
                    case "BASIS OF ENCOUNTER":
                        boe = getDF(dataFrames, 'BASIS OF ENCOUNTER')
                        fillTable(boe, table, startrow, startcol)

def fillPresentation(templatePath, filePath, summaryDataFrames, breakdownDataFrames, dateStart, dateEnd):
    presentation = Presentation(templatePath)
    summaryPage(presentation, summaryDataFrames, dateStart, dateEnd)
    breakdownPages(presentation, breakdownDataFrames, dateStart, dateEnd)
    presentation.save(filePath)